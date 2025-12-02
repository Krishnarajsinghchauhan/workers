import json
import sys
import boto3
import tempfile
import os
from pdf2docx import Converter
import camelot
from pptx import Presentation
from pdf2image import convert_from_path

s3 = boto3.client("s3")
BUCKET = "pdf-master-storage"


def pdf_to_word(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()


def pdf_to_excel(input_path, output_path):
    tables = camelot.read_pdf(input_path, pages="all")
    tables.export(output_path, f="excel")


def pdf_to_ppt(input_path, output_path):
    pages = convert_from_path(input_path)
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for img in pages:
        slide = prs.slides.add_slide(blank)
        temp_img = tempfile.mktemp(suffix=".png")
        img.save(temp_img)
        slide.shapes.add_picture(temp_img, 0, 0, prs.slide_width, prs.slide_height)

    prs.save(output_path)


def download_s3_file(url):
    key = url.replace(f"https://{BUCKET}.s3.amazonaws.com/", "")
    temp = tempfile.mktemp(suffix=".pdf")
    s3.download_file(BUCKET, key, temp)
    return temp


def upload_file(path, job_id):
    ext = os.path.splitext(path)[1]
    key = f"processed/{job_id}{ext}"
    s3.upload_file(path, BUCKET, key)
    return f"https://{BUCKET}.s3.amazonaws.com/{key}"


def main():
    try:
        job = json.loads(sys.stdin.read())

        tool = job["tool"]
        job_id = job["job_id"]
        file_url = job["files"][0]

        input_path = download_s3_file(file_url)

        if tool == "pdf-to-word":
            out = tempfile.mktemp(suffix=".docx")
            pdf_to_word(input_path, out)

        elif tool == "pdf-to-excel":
            out = tempfile.mktemp(suffix=".xlsx")
            pdf_to_excel(input_path, out)

        elif tool == "pdf-to-ppt":
            out = tempfile.mktemp(suffix=".pptx")
            pdf_to_ppt(input_path, out)

        url = upload_file(out, job_id)

        print(json.dumps({"status": "completed", "url": url}))

    except Exception as e:
        print(json.dumps({
            "status": "error",
            "error": str(e)
        }))
        sys.exit(1)


if __name__ == "__main__":
    main()
